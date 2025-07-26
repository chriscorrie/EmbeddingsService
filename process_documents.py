import os
import logging
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

# Legacy Office format support
try:
    import docx2txt
    DOCX2TXT_AVAILABLE = True
except ImportError:
    DOCX2TXT_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

logger = logging.getLogger(__name__)

def extract_text_from_file(file_path):
    """
    Extract text from various file formats including legacy Office formats
    
    Args:
        file_path (str): Path to the file
        
    Returns:
        str: Extracted text content, or None if extraction fails or format is unsupported
    """
    try:
        file_lower = file_path.lower()
        
        if file_lower.endswith('.pdf'):
            return extract_text_from_pdf(file_path)
        elif file_lower.endswith('.docx'):
            return extract_text_from_docx(file_path)
        elif file_lower.endswith('.doc'):
            return extract_text_from_legacy_doc(file_path)
        elif file_lower.endswith('.xlsx'):
            return extract_text_from_excel(file_path)
        elif file_lower.endswith('.xls'):
            return extract_text_from_legacy_excel(file_path)
        elif file_lower.endswith('.pptx'):
            return extract_text_from_pptx(file_path)
        elif file_lower.endswith('.ppt'):
            return extract_text_from_legacy_ppt(file_path)
        elif file_lower.endswith('.txt'):
            return extract_text_from_txt(file_path)
        elif file_lower.endswith('.rtf'):
            return extract_text_from_rtf(file_path)
        else:
            # Return None for unsupported file types instead of raising exception
            logger.debug(f"Unsupported file format: {file_path}")
            return None
    except Exception as e:
        # Return None for any extraction errors to allow graceful handling
        logger.warning(f"Error extracting text from {file_path}: {e}")
        return None

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    return '\n'.join(page.extract_text() for page in reader.pages)

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return '\n'.join(paragraph.text for paragraph in doc.paragraphs)

def extract_text_from_excel(file_path):
    workbook = load_workbook(file_path, data_only=True)
    text = []
    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        for row in worksheet.iter_rows(values_only=True):
            text.append(' '.join(map(str, row)))
    return '\n'.join(text)

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text.append(shape.text)
    return '\n'.join(text)

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def extract_text_from_legacy_doc(file_path):
    """
    Extract text from legacy .doc files using docx2txt
    
    Args:
        file_path (str): Path to the .doc file
        
    Returns:
        str: Extracted text content, or None if extraction fails
    """
    if not DOCX2TXT_AVAILABLE:
        logger.warning("docx2txt not available - cannot process .doc files")
        return None
    
    try:
        # docx2txt can handle both .docx and .doc files
        text_content = docx2txt.process(file_path)
        return text_content if text_content else None
    except Exception as e:
        logger.warning(f"Failed to extract text from .doc file {file_path}: {e}")
        return None

def extract_text_from_legacy_excel(file_path):
    """
    Extract text from legacy .xls files using xlrd
    
    Args:
        file_path (str): Path to the .xls file
        
    Returns:
        str: Extracted text content, or None if extraction fails
    """
    if not XLRD_AVAILABLE:
        logger.warning("xlrd not available - cannot process .xls files")
        return None
    
    try:
        workbook = xlrd.open_workbook(file_path)
        text = []
        
        for sheet_idx in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_idx)
            for row_idx in range(sheet.nrows):
                row_values = []
                for col_idx in range(sheet.ncols):
                    cell_value = sheet.cell_value(row_idx, col_idx)
                    if cell_value:
                        row_values.append(str(cell_value))
                if row_values:
                    text.append(' '.join(row_values))
        
        return '\n'.join(text)
    except Exception as e:
        logger.warning(f"Failed to extract text from .xls file {file_path}: {e}")
        return None

def extract_text_from_legacy_ppt(file_path):
    """
    Extract text from legacy .ppt files using LibreOffice conversion
    
    Args:
        file_path (str): Path to the .ppt file
        
    Returns:
        str: Extracted text content, or None if extraction fails
    """
    # For now, we'll skip .ppt files as they require LibreOffice conversion
    # which is more complex to set up. We can implement this later if needed.
    logger.warning(f"Legacy .ppt file format not yet supported: {file_path}")
    return None

def extract_text_from_rtf(file_path):
    """
    Extract text from RTF files - basic RTF support
    
    Args:
        file_path (str): Path to the .rtf file
        
    Returns:
        str: Extracted text content, or None if extraction fails
    """
    try:
        # Simple RTF text extraction - read as text and strip RTF formatting
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            content = file.read()
            # Very basic RTF parsing - extract text between spaces
            import re
            # Remove RTF control sequences
            text = re.sub(r'\\[a-z]+\d*\s?', ' ', content)
            text = re.sub(r'[{}]', '', text)
            # Clean up extra whitespace
            text = re.sub(r'\s+', ' ', text).strip()
            return text if text else None
    except Exception as e:
        logger.warning(f"Failed to extract text from .rtf file {file_path}: {e}")
        return None

if __name__ == '__main__':
    # Example usage
    sample_file = '/mnt/HomerShare/FBO Attachments/sample.pdf'
    print(extract_text_from_file(sample_file))
